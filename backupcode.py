import streamlit as st
import os, io, re, shutil, zipfile
from io import BytesIO
from typing import List, Dict
import numpy as np

# # Fix for ChromaDB on Streamlit Cloud
# import sys
# __import__('pysqlite3')
# sys.modules['sqlite3'] = sys.modules.pop('pysqlite3')

# ----------- CORE LIBS -----------
import PyPDF2
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer, CrossEncoder
import requests

# ----------- DOCUMENT PROCESSING -----------
try:
    from docx import Document
    DOCX_AVAILABLE = True
except:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except:
    PPTX_AVAILABLE = False

import markdown
from bs4 import BeautifulSoup

# ================= STREAMLIT CONFIG =================
st.set_page_config(
    page_title="📚 Advanced RAG Assistant",
    page_icon="📚",
    layout="wide"
)

st.title("📚 Advanced RAG Syllabus Assistant")
st.caption("Multi-stage retrieval with semantic chunking and cross-encoder reranking")

# ================= SECRETS =================
HF_TOKEN = st.secrets["HF_TOKEN"]
HF_API_URL = "https://router.huggingface.co/v1/chat/completions"
MODEL_NAME = "meta-llama/Llama-3.1-8B-Instruct:fastest"
HEADERS = {
    "Authorization": f"Bearer {HF_TOKEN}",
    "Content-Type": "application/json"
}

# ================= CACHE MODELS =================
@st.cache_resource
def load_models():
    """Load embedding model, cross-encoder for reranking, and ChromaDB client"""
    embed = SentenceTransformer("all-MiniLM-L6-v2")
    reranker = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2')
    client = chromadb.Client(
        Settings(
            persist_directory="./chroma_db",
            anonymized_telemetry=False
        )
    )
    return embed, reranker, client

embedding_model, cross_encoder, chroma_client = load_models()

# ================= CONSTANTS =================
# Chunking parameters
SEMANTIC_CHUNK_SIZE = 400
SLIDING_WINDOW_SIZE = 600
WINDOW_OVERLAP = 150

# Retrieval parameters
INITIAL_RETRIEVAL_K = 20
RERANK_TOP_K = 10
FINAL_CONTEXT_K = 5

HEADER_FOOTER_PATTERNS = [
    r'^\d+$',
    r'^Page\s+\d+',
    r'©.*\d{4}',
    r'^Chapter\s+\d+$',
    r'^Section\s+\d+$',
    r'^\d+\s*/\s*\d+$',
]

SUPPORTED_EXTENSIONS = ['.pdf', '.txt', '.md']
if DOCX_AVAILABLE:
    SUPPORTED_EXTENSIONS.extend(['.docx', '.doc'])
if PPTX_AVAILABLE:
    SUPPORTED_EXTENSIONS.extend(['.pptx', '.ppt'])

# ================= HEADER/FOOTER REMOVAL =================
def is_header_footer(text: str) -> bool:
    text = text.strip()
    if len(text) < 3 or len(text) > 100:
        return len(text) < 3
    
    for pattern in HEADER_FOOTER_PATTERNS:
        if re.match(pattern, text, re.IGNORECASE):
            return True
    return False

def clean_text(text: str) -> str:
    lines = text.split('\n')
    cleaned = [line for line in lines if not is_header_footer(line)]
    result = '\n'.join(cleaned)
    result = re.sub(r'\n{3,}', '\n\n', result)
    return result.strip()

# ================= ZIP FILE HANDLING =================
def extract_files_from_zip(zip_bytes: bytes) -> List[Dict]:
    """Extract all supported files from a ZIP archive."""
    extracted_files = []
    
    try:
        with zipfile.ZipFile(BytesIO(zip_bytes), 'r') as zip_ref:
            for file_info in zip_ref.filelist:
                if file_info.is_dir():
                    continue
                
                filename = file_info.filename
                ext = os.path.splitext(filename)[1].lower()
                
                if ext in SUPPORTED_EXTENSIONS:
                    try:
                        file_bytes = zip_ref.read(file_info)
                        clean_name = os.path.basename(filename)
                        extracted_files.append({
                            'name': clean_name,
                            'bytes': file_bytes
                        })
                    except Exception as e:
                        st.warning(f"⚠️ Could not extract {filename}: {str(e)}")
    except Exception as e:
        st.error(f"❌ Error reading ZIP file: {str(e)}")
    
    return extracted_files

# ================= TEXT EXTRACTION BY FILE TYPE =================
def extract_text_from_txt(file_bytes: bytes) -> List[Dict]:
    try:
        text = file_bytes.decode('utf-8')
    except:
        text = file_bytes.decode('latin-1')
    text = clean_text(text)
    return [{"page": 1, "text": text}] if text else []

def extract_text_from_md(file_bytes: bytes) -> List[Dict]:
    try:
        md_text = file_bytes.decode('utf-8')
    except:
        md_text = file_bytes.decode('latin-1')
    html = markdown.markdown(md_text)
    soup = BeautifulSoup(html, 'html.parser')
    text = soup.get_text()
    text = clean_text(text)
    return [{"page": 1, "text": text}] if text else []

def extract_text_from_docx(file_bytes: bytes) -> List[Dict]:
    if not DOCX_AVAILABLE:
        st.error("DOCX support not available")
        return []
    doc = Document(BytesIO(file_bytes))
    full_text = [para.text for para in doc.paragraphs if para.text.strip()]
    text = '\n'.join(full_text)
    text = clean_text(text)
    return [{"page": 1, "text": text}] if text else []

def extract_text_from_pptx(file_bytes: bytes) -> List[Dict]:
    if not PPTX_AVAILABLE:
        st.error("PPTX support not available")
        return []
    prs = Presentation(BytesIO(file_bytes))
    pages = []
    
    for i, slide in enumerate(prs.slides):
        text_parts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
        text = '\n'.join(text_parts)
        text = clean_text(text)
        if text:
            pages.append({"page": i + 1, "text": text})
    return pages

def extract_text_from_pdf(pdf_bytes: bytes) -> List[Dict]:
    """Extract text from PDF using PyPDF2 (basic extraction, no OCR)"""
    reader = PyPDF2.PdfReader(io.BytesIO(pdf_bytes))
    pages = []
    
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
            text = text.strip()
            text = clean_text(text)
            
            if text:
                pages.append({"page": i + 1, "text": text})
        except Exception as e:
            st.warning(f"Could not extract page {i+1}: {str(e)}")
    
    return pages

def extract_text_from_file(file_bytes: bytes, filename: str) -> List[Dict]:
    ext = os.path.splitext(filename)[1].lower()
    extractors = {
        '.pdf': extract_text_from_pdf,
        '.txt': extract_text_from_txt,
        '.md': extract_text_from_md,
    }
    
    if DOCX_AVAILABLE:
        extractors['.docx'] = extract_text_from_docx
        extractors['.doc'] = extract_text_from_docx
    
    if PPTX_AVAILABLE:
        extractors['.pptx'] = extract_text_from_pptx
        extractors['.ppt'] = extract_text_from_pptx
    
    extractor = extractors.get(ext)
    if extractor:
        try:
            return extractor(file_bytes)
        except Exception as e:
            st.error(f"Error extracting {filename}: {str(e)}")
            return []
    else:
        st.warning(f"Unsupported file type: {ext}")
        return []

# ================= ADVANCED CHUNKING =================
def semantic_chunking(text: str, chunk_size: int = SEMANTIC_CHUNK_SIZE) -> List[str]:
    """Semantic chunking: splits text at sentence boundaries."""
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = []
    current_size = 0
    
    for sentence in sentences:
        words = sentence.split()
        sentence_size = len(words)
        
        if current_size + sentence_size > chunk_size and current_chunk:
            chunks.append(' '.join(current_chunk))
            current_chunk = [sentence]
            current_size = sentence_size
        else:
            current_chunk.append(sentence)
            current_size += sentence_size
    
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks

def sliding_window_chunking(text: str, window_size: int = SLIDING_WINDOW_SIZE, 
                           overlap: int = WINDOW_OVERLAP) -> List[str]:
    """Sliding window chunking: creates overlapping chunks."""
    words = text.split()
    chunks = []
    
    for i in range(0, len(words), window_size - overlap):
        chunk = ' '.join(words[i:i + window_size])
        if chunk.strip():
            chunks.append(chunk)
        
        if i + window_size >= len(words):
            break
    
    return chunks

def hybrid_chunking(pages: List[Dict]) -> List[Dict]:
    """Hybrid approach: combines semantic and sliding window chunking."""
    all_chunks = []
    chunk_id = 0
    
    for p in pages:
        text = p["text"]
        
        semantic_chunks = semantic_chunking(text, SEMANTIC_CHUNK_SIZE)
        
        for sem_chunk in semantic_chunks:
            all_chunks.append({
                "page": p["page"],
                "text": sem_chunk,
                "doc": p.get("doc", "unknown"),
                "chunk_type": "semantic",
                "chunk_id": chunk_id
            })
            chunk_id += 1
        
        if len(text.split()) > SLIDING_WINDOW_SIZE:
            window_chunks = sliding_window_chunking(text, SLIDING_WINDOW_SIZE, WINDOW_OVERLAP)
            
            for win_chunk in window_chunks:
                all_chunks.append({
                    "page": p["page"],
                    "text": win_chunk,
                    "doc": p.get("doc", "unknown"),
                    "chunk_type": "window",
                    "chunk_id": chunk_id
                })
                chunk_id += 1
    
    return all_chunks

# ================= MULTI-STAGE RETRIEVAL WITH RERANKING =================
def retrieve_and_rerank(question: str, collection, top_k: int = FINAL_CONTEXT_K) -> List[Dict]:
    """Multi-stage retrieval with reranking and deduplication."""
    q_emb = embedding_model.encode([question]).tolist()
    res = collection.query(q_emb, n_results=INITIAL_RETRIEVAL_K)
    
    candidates = []
    for doc, metadata, distance in zip(res["documents"][0], res["metadatas"][0], res["distances"][0]):
        candidates.append({
            "text": doc,
            "page": metadata.get("page", 1),
            "doc": metadata.get("doc", "unknown"),
            "chunk_type": metadata.get("chunk_type", "unknown"),
            "initial_score": 1 - distance
        })
    
    if not candidates:
        return []
    
    pairs = [[question, c["text"]] for c in candidates]
    rerank_scores = cross_encoder.predict(pairs)
    
    for i, score in enumerate(rerank_scores):
        candidates[i]["rerank_score"] = float(score)
    
    candidates.sort(key=lambda x: x["rerank_score"], reverse=True)
    
    final_chunks = []
    seen_texts = []
    
    for candidate in candidates[:RERANK_TOP_K]:
        is_duplicate = False
        for seen in seen_texts:
            overlap = len(set(candidate["text"].split()) & set(seen.split()))
            total = len(set(candidate["text"].split()) | set(seen.split()))
            if total > 0 and overlap / total > 0.7:
                is_duplicate = True
                break
        
        if not is_duplicate:
            final_chunks.append(candidate)
            seen_texts.append(candidate["text"])
            
            if len(final_chunks) >= top_k:
                break
    
    return final_chunks

# ================= MULTI-QUESTION PARSING =================
def parse_multiple_questions(user_input: str) -> List[str]:
    """Parse multiple questions from user input"""
    questions = re.split(r'\n+|\?+(?=\s*[A-Z0-9])', user_input)
    parsed = []
    
    for q in questions:
        q = q.strip()
        if q and len(q) > 5:
            if not q.endswith('?'):
                q += '?'
            parsed.append(q)
    
    if not parsed:
        parsed = [user_input.strip()]
    
    return parsed

# ================= ENHANCED LLM QUERY =================
def llm_query(question: str, chunks: List[Dict]) -> str:
    """Enhanced LLM query with structured prompting."""
    if not chunks:
        return "❌ Answer not found in the documents."
    
    context_parts = []
    for i, c in enumerate(chunks, 1):
        score_info = f"[Relevance: {c.get('rerank_score', 0):.2f}]"
        source_info = f"**Source {i}** (Page {c['page']}, {c.get('doc', 'document')}) {score_info}:"
        context_parts.append(f"{source_info}\n{c['text']}\n")
    
    context = "\n".join(context_parts)
    
    prompt = f"""You are a knowledgeable teaching assistant. Answer the question using ONLY the document content provided below.

**INSTRUCTIONS:**
1. Provide a comprehensive, well-structured answer
2. Use bullet points for listing multiple items or key points
3. Use paragraphs for explanations and descriptions
4. Include specific details, dates, numbers, and examples from the documents
5. If the answer involves multiple aspects, organize them clearly
6. Add context and brief explanations to help understanding
7. Cite the source number (e.g., "Source 1") when referencing specific information
8. If the answer is not in the documents, clearly state: "Answer not found in the documents."

**DOCUMENT CONTENT:**
{context}

**QUESTION:** {question}

**ANSWER:**"""
    
    payload = {
        "model": MODEL_NAME,
        "messages": [
            {
                "role": "system", 
                "content": "You are an expert teaching assistant who provides clear, well-structured, and comprehensive answers based on document content. You format answers with appropriate structure including bullet points, paragraphs, and explanations for easy understanding."
            },
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.2,
        "max_tokens": 600,
        "top_p": 0.9
    }
    
    try:
        r = requests.post(HF_API_URL, headers=HEADERS, json=payload, timeout=30)
        
        if r.status_code != 200:
            return f"⚠️ API Error: {r.status_code}"
        
        result = r.json()
        answer = result["choices"][0]["message"]["content"].strip()
        
        return answer
    except Exception as e:
        return f"⚠️ Error querying LLM: {str(e)}"

def process_multiple_questions(questions: List[str], collection) -> Dict[str, Dict]:
    """Process multiple questions with enhanced retrieval"""
    results = {}
    
    for q in questions:
        chunks = retrieve_and_rerank(q, collection, FINAL_CONTEXT_K)
        answer = llm_query(q, chunks)
        
        results[q] = {
            "answer": answer,
            "chunks_used": len(chunks),
            "sources": [
                {
                    "doc": c.get("doc"), 
                    "page": c.get("page"), 
                    "score": c.get("rerank_score", 0)
                } for c in chunks
            ]
        }
    
    return results

# ================= SESSION STATE =================
if "ready" not in st.session_state:
    st.session_state.ready = False
if "processed_files" not in st.session_state:
    st.session_state.processed_files = []
if "chunk_stats" not in st.session_state:
    st.session_state.chunk_stats = {}
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

# ================= UI =================
st.header("📤 Upload Documents")

# Display supported formats
supported_str = "PDF, TXT, MD"
if DOCX_AVAILABLE:
    supported_str += ", DOCX"
if PPTX_AVAILABLE:
    supported_str += ", PPTX"
st.caption(f"Supported formats: {supported_str}, ZIP")

# Settings in sidebar
with st.sidebar:
    st.header("⚙️ Advanced Settings")
    
    st.subheader("Retrieval Parameters")
    initial_k = st.slider("Initial Retrieval K", 10, 50, INITIAL_RETRIEVAL_K, 5,
                         help="Number of chunks to retrieve initially")
    final_k = st.slider("Final Context K", 3, 15, FINAL_CONTEXT_K, 1,
                       help="Number of top chunks to use for answer")
    
    st.subheader("Chunking Parameters")
    semantic_size = st.slider("Semantic Chunk Size (words)", 200, 800, SEMANTIC_CHUNK_SIZE, 50)
    window_size = st.slider("Sliding Window Size (words)", 300, 1000, SLIDING_WINDOW_SIZE, 50)
    window_overlap = st.slider("Window Overlap (words)", 50, 300, WINDOW_OVERLAP, 25)

file_types = ["pdf", "txt", "md", "zip"]
if DOCX_AVAILABLE:
    file_types.extend(["docx", "doc"])
if PPTX_AVAILABLE:
    file_types.extend(["pptx", "ppt"])

uploaded_files = st.file_uploader(
    "Upload documents or ZIP files containing documents",
    type=file_types,
    accept_multiple_files=True
)

if uploaded_files and st.button("🚀 Process All Documents with Advanced Chunking"):
    with st.spinner("Processing documents with hybrid chunking strategy..."):
        all_pages = []
        files_to_process = []
        
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        status_text.info("📦 Extracting files from ZIP archives...")
        for uploaded in uploaded_files:
            if uploaded.name.lower().endswith('.zip'):
                zip_bytes = uploaded.read()
                extracted = extract_files_from_zip(zip_bytes)
                files_to_process.extend(extracted)
                st.success(f"✅ Extracted {len(extracted)} files from {uploaded.name}")
            else:
                files_to_process.append({
                    'name': uploaded.name,
                    'bytes': uploaded.read()
                })
        
        if not files_to_process:
            st.error("❌ No valid files found to process")
            st.stop()
        
        st.info(f"📄 Processing {len(files_to_process)} total files...")
        
        for idx, file_data in enumerate(files_to_process):
            status_text.info(f"📄 Processing: {file_data['name']}")
            
            pages = extract_text_from_file(file_data['bytes'], file_data['name'])
            
            for p in pages:
                p["doc"] = file_data['name']
            
            all_pages.extend(pages)
            progress_bar.progress((idx + 1) / len(files_to_process))
        
        status_text.info("✂️ Applying hybrid chunking (semantic + sliding window)...")
        chunks = hybrid_chunking(all_pages)
        
        if chunks:
            texts = [c["text"] for c in chunks]
            
            status_text.info("🧠 Generating embeddings...")
            embeddings = embedding_model.encode(texts, show_progress_bar=False).tolist()
            
            try:
                chroma_client.delete_collection("syllabus")
            except:
                pass
            
            status_text.info("💾 Storing in vector database...")
            col = chroma_client.create_collection("syllabus")
            col.add(
                documents=texts,
                embeddings=embeddings,
                metadatas=[
                    {
                        "page": c["page"], 
                        "doc": c.get("doc", "unknown"),
                        "chunk_type": c.get("chunk_type", "unknown")
                    } for c in chunks
                ],
                ids=[f"c{i}" for i in range(len(texts))]
            )
            
            semantic_count = sum(1 for c in chunks if c.get("chunk_type") == "semantic")
            window_count = sum(1 for c in chunks if c.get("chunk_type") == "window")
            
            st.session_state.chunk_stats = {
                "total": len(chunks),
                "semantic": semantic_count,
                "window": window_count
            }
            st.session_state.ready = True
            st.session_state.processed_files = [f['name'] for f in files_to_process]
            st.session_state.chat_history = []
            
            status_text.empty()
            progress_bar.empty()
            
            st.success("✅ Processing Complete!")
            
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Files Processed", len(files_to_process))
            with col2:
                st.metric("Total Chunks", len(chunks))
            with col3:
                st.metric("Semantic Chunks", semantic_count)
            with col4:
                st.metric("Window Chunks", window_count)
            
            with st.expander("📄 Processed Files", expanded=False):
                for fname in st.session_state.processed_files:
                    st.text(f"✓ {fname}")
        else:
            st.error("❌ No content extracted from documents")

# ================= QUERY INTERFACE =================
if st.session_state.ready:
    st.divider()
    st.header("💬 Ask Questions")
    
    with st.expander("📊 Document Statistics", expanded=False):
        stats = st.session_state.chunk_stats
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Total Chunks", stats.get("total", 0))
        with col2:
            st.metric("Semantic Chunks", stats.get("semantic", 0))
        with col3:
            st.metric("Sliding Window Chunks", stats.get("window", 0))
    
    st.info("💡 **Tip:** You can ask single or multiple questions. Separate multiple questions with new lines for batch processing.")
    
    # Display chat history
    for i, chat in enumerate(st.session_state.chat_history):
        with st.container():
            if isinstance(chat['question'], list):
                st.markdown("**👤 You asked:**")
                for idx, q in enumerate(chat['question'], 1):
                    st.markdown(f"{idx}. {q}")
            else:
                st.markdown(f"**👤 You:** {chat['question']}")
            
            st.markdown("")
            
            if isinstance(chat['answer'], dict):
                for idx, (question, answer_data) in enumerate(chat['answer'].items(), 1):
                    st.markdown(f"**🤖 Assistant (Question {idx}):**")
                    st.markdown(f"""
                    <div style="background-color: #000000; padding: 20px; border-radius: 10px; border-left: 5px solid #4CAF50; margin-bottom: 15px;">
                    <strong style="color: #1f77b4;">Q{idx}: {question}</strong><br><br>
                    {answer_data['answer']}
                    </div>
                    """, unsafe_allow_html=True)
                    
                    with st.expander(f"📚 Sources for Question {idx}", expanded=False):
                        st.caption(f"**Chunks analyzed:** {answer_data.get('chunks_used', 0)}")
                        for j, source in enumerate(answer_data.get('sources', []), 1):
                            score_color = "#4CAF50" if source['score'] > 0.5 else "#FF9800"
                            st.markdown(f"""
                            **Source {j}:** {source['doc']} (Page {source['page']})  
                            <span style="color: {score_color}">Relevance Score: {source['score']:.3f}</span>
                            """, unsafe_allow_html=True)
            else:
                st.markdown("**🤖 Assistant:**")
                st.markdown(f"""
                <div style="background-color: #000000; padding: 20px; border-radius: 10px; border-left: 5px solid #4CAF50; margin-bottom: 10px;">
                {chat['answer']}
                </div>
                """, unsafe_allow_html=True)
                
                with st.expander(f"📚 Sources", expanded=False):
                    st.caption(f"**Chunks analyzed:** {chat.get('chunks_used', 0)}")
                    for j, source in enumerate(chat.get('sources', []), 1):
                        score_color = "#4CAF50" if source['score'] > 0.5 else "#FF9800"
                        st.markdown(f"""
                        **Source {j}:** {source['doc']} (Page {source['page']})  
                        <span style="color: {score_color}">Relevance Score: {source['score']:.3f}</span>
                        """, unsafe_allow_html=True)
            
            st.divider()
    
    user_input = st.chat_input("Ask one or multiple questions (separate with new lines)...")
    
    if user_input:
        col = chroma_client.get_collection("syllabus")
        
        questions = parse_multiple_questions(user_input)
        
        if len(questions) > 1:
            with st.spinner(f"🔄 Processing {len(questions)} questions..."):
                results = process_multiple_questions(questions, col)
                
                chat_entry = {
                    "question": questions,
                    "answer": results,
                    "is_batch": True
                }
                st.session_state.chat_history.append(chat_entry)
        else:
            with st.spinner("🔄 Retrieving and analyzing..."):
                chunks = retrieve_and_rerank(questions[0], col, final_k)
                answer = llm_query(questions[0], chunks)
                
                chat_entry = {
                    "question": questions[0],
                    "answer": answer,
                    "chunks_used": len(chunks),
                    "sources": [
                        {
                            "doc": c.get("doc"), 
                            "page": c.get("page"), 
                            "score": c.get("rerank_score", 0)
                        } for c in chunks
                    ],
                    "is_batch": False
                }
                st.session_state.chat_history.append(chat_entry)
        
        st.rerun()

# ================= CONTROLS =================
st.divider()
col1, col2, col3, col4 = st.columns(4)

with col1:
    if st.button("🔄 Reset Database"):
        try:
            chroma_client.delete_collection("syllabus")
        except:
            pass
        if os.path.exists("./chroma_db"):
            shutil.rmtree("./chroma_db", ignore_errors=True)
        st.session_state.ready = False
        st.session_state.processed_files = []
        st.session_state.chunk_stats = {}
        st.session_state.chat_history = []
        st.success("✅ Reset complete")
        st.rerun()

with col2:
    if st.session_state.ready and st.button("🗑️ Clear Chat History"):
        st.session_state.chat_history = []
        st.success("✅ Chat cleared")
        st.rerun()

with col3:
    if st.session_state.ready:
        st.metric("📚 Documents", len(st.session_state.processed_files))

with col4:
    if st.session_state.ready:
        st.metric("💬 Messages", len(st.session_state.chat_history))

# ================= FOOTER =================
st.divider()
st.caption("""
**Advanced Features:**
- 🎯 Multi-stage retrieval with cross-encoder reranking
- ✂️ Hybrid chunking: semantic + sliding window
- 🔍 Deduplication to avoid redundant content
- 📊 Relevance scoring for transparency
- 🎨 Enhanced answer formatting
- 💬 Chat mode with single or batch question support
- 📋 Intelligent parsing: automatically detects multiple questions
- 📦 ZIP file support for bulk uploads
- ⚡ Optimized for Streamlit Cloud deployment
""")