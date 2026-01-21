import streamlit as st
import os, io, math, re, shutil
from io import BytesIO
from typing import List, Dict, Tuple
from collections import defaultdict
import numpy as np

# ----------- CORE LIBS -----------
import PyPDF2
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer, CrossEncoder
import requests

# ----------- OCR & IMAGE -----------
from PIL import Image, ImageEnhance
import pytesseract
from pdf2image import convert_from_bytes
import fitz  # PyMuPDF
import imagehash

# ----------- DOCUMENT PROCESSING -----------
from docx import Document  # python-docx
from pptx import Presentation  # python-pptx
import markdown
from bs4 import BeautifulSoup

# ================= STREAMLIT CONFIG =================
st.set_page_config(
    page_title="📚 Advanced RAG Assistant",
    page_icon="📚",
    layout="wide"
)

st.title("📚 Advanced RAG Syllabus Assistant")
st.caption("Multi-stage retrieval with semantic chunking and cross-encoder reranking")

# ================= SECRETS =================
# HF_TOKEN = st.secrets["HF_TOKEN"]
HF_TOKEN = st.secrets["HF_TOKEN"]
HF_API_URL = "https://router.huggingface.co/v1/chat/completions"
MODEL_NAME = "meta-llama/Llama-3.1-8B-Instruct:fastest"
HEADERS = {
    "Authorization": f"Bearer {HF_TOKEN}",
    "Content-Type": "application/json"
}

# ================= CACHE MODELS =================
@st.cache_resource
def load_models():
    """Load embedding model, cross-encoder for reranking, and ChromaDB client"""
    embed = SentenceTransformer("all-MiniLM-L6-v2")
    # Cross-encoder for reranking retrieved chunks
    reranker = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2')
    client = chromadb.Client(
        Settings(
            persist_directory="./chroma_db",
            anonymized_telemetry=False
        )
    )
    return embed, reranker, client

embedding_model, cross_encoder, chroma_client = load_models()

# ================= CONSTANTS =================
OCR_THRESHOLD = 50
IMAGE_DIR = "data/images"
os.makedirs(IMAGE_DIR, exist_ok=True)

# Chunking parameters
SEMANTIC_CHUNK_SIZE = 400  # Words per semantic chunk
SLIDING_WINDOW_SIZE = 600  # Words for sliding window
WINDOW_OVERLAP = 150  # Overlap between windows

# Retrieval parameters
INITIAL_RETRIEVAL_K = 20  # Retrieve more initially
RERANK_TOP_K = 10  # Keep top K after reranking
FINAL_CONTEXT_K = 5  # Use top K for final answer

HEADER_FOOTER_PATTERNS = [
    r'^\d+$',
    r'^Page\s+\d+',
    r'©.*\d{4}',
    r'^Chapter\s+\d+$',
    r'^Section\s+\d+$',
    r'^\d+\s*/\s*\d+$',
]

# ================= HEADER/FOOTER REMOVAL =================
def is_header_footer(text: str) -> bool:
    text = text.strip()
    if len(text) < 3 or len(text) > 100:
        return len(text) < 3
    
    for pattern in HEADER_FOOTER_PATTERNS:
        if re.match(pattern, text, re.IGNORECASE):
            return True
    return False

def clean_text(text: str) -> str:
    lines = text.split('\n')
    cleaned = [line for line in lines if not is_header_footer(line)]
    result = '\n'.join(cleaned)
    result = re.sub(r'\n{3,}', '\n\n', result)
    return result.strip()

# ================= TEXT EXTRACTION BY FILE TYPE =================
def extract_text_from_txt(file_bytes: bytes) -> List[Dict]:
    try:
        text = file_bytes.decode('utf-8')
    except:
        text = file_bytes.decode('latin-1')
    text = clean_text(text)
    return [{"page": 1, "text": text}] if text else []

def extract_text_from_md(file_bytes: bytes) -> List[Dict]:
    try:
        md_text = file_bytes.decode('utf-8')
    except:
        md_text = file_bytes.decode('latin-1')
    html = markdown.markdown(md_text)
    soup = BeautifulSoup(html, 'html.parser')
    text = soup.get_text()
    text = clean_text(text)
    return [{"page": 1, "text": text}] if text else []

def extract_text_from_docx(file_bytes: bytes) -> List[Dict]:
    doc = Document(BytesIO(file_bytes))
    full_text = [para.text for para in doc.paragraphs if para.text.strip()]
    text = '\n'.join(full_text)
    text = clean_text(text)
    return [{"page": 1, "text": text}] if text else []

def extract_text_from_pptx(file_bytes: bytes) -> List[Dict]:
    prs = Presentation(BytesIO(file_bytes))
    pages = []
    
    for i, slide in enumerate(prs.slides):
        text_parts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
        text = '\n'.join(text_parts)
        text = clean_text(text)
        if text:
            pages.append({"page": i + 1, "text": text})
    return pages

def extract_text_from_pdf(pdf_bytes: bytes) -> List[Dict]:
    reader = PyPDF2.PdfReader(io.BytesIO(pdf_bytes))
    pages = []
    
    try:
        images = convert_from_bytes(pdf_bytes, dpi=200)
    except:
        images = []
    
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        text = text.strip()
        
        if len(text) < OCR_THRESHOLD and i < len(images):
            img = images[i].convert("L")
            ocr_text = pytesseract.image_to_string(img)
            if len(ocr_text.strip()) > len(text):
                text = ocr_text.strip()
        
        text = clean_text(text)
        if text:
            pages.append({"page": i + 1, "text": text})
    
    return pages

def extract_text_from_file(file_bytes: bytes, filename: str) -> List[Dict]:
    ext = os.path.splitext(filename)[1].lower()
    extractors = {
        '.pdf': extract_text_from_pdf,
        '.txt': extract_text_from_txt,
        '.md': extract_text_from_md,
        '.docx': extract_text_from_docx,
        '.doc': extract_text_from_docx,
        '.pptx': extract_text_from_pptx,
        '.ppt': extract_text_from_pptx,
    }
    
    extractor = extractors.get(ext)
    if extractor:
        try:
            return extractor(file_bytes)
        except Exception as e:
            st.error(f"Error extracting {filename}: {str(e)}")
            return []
    else:
        st.warning(f"Unsupported file type: {ext}")
        return []

# ================= IMAGE UTILITIES =================
def image_entropy(img: Image.Image) -> float:
    hist = img.histogram()
    total = sum(hist)
    entropy = 0
    for c in hist:
        if c == 0:
            continue
        p = c / total
        entropy -= p * math.log2(p)
    return entropy

def extract_images_from_pdf(pdf_bytes: bytes, doc_name: str) -> List[Dict]:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    seen_hashes = []
    image_data = []
    
    for page_index in range(len(doc)):
        page = doc[page_index]
        images = page.get_images(full=True)
        
        for img in images:
            base = doc.extract_image(img[0])
            img_bytes = base["image"]
            ext = base["ext"]
            
            pil = Image.open(BytesIO(img_bytes)).convert("RGB")
            
            if image_entropy(pil) < 4.0:
                continue
            
            h = imagehash.phash(pil)
            if any(abs(h - sh) <= 5 for sh in seen_hashes):
                continue
            
            seen_hashes.append(h)
            safe_name = re.sub(r'[^\w\-_.]', '_', doc_name)
            path = os.path.join(IMAGE_DIR, f"{safe_name}_page{page_index+1}_{len(seen_hashes)}.{ext}")
            
            with open(path, "wb") as f:
                f.write(img_bytes)
            
            image_data.append({"page": page_index + 1, "path": path, "doc": doc_name})
    
    doc.close()
    return image_data

def ocr_image_text(path: str) -> str:
    try:
        img = Image.open(path).convert("L")
        img = ImageEnhance.Contrast(img).enhance(2.0)
        text = pytesseract.image_to_string(img, config="--psm 6")
        return re.sub(r"[^a-zA-Z0-9\s\.,:-]", "", text).strip()
    except:
        return ""

# ================= ADVANCED CHUNKING =================
def semantic_chunking(text: str, chunk_size: int = SEMANTIC_CHUNK_SIZE) -> List[str]:
    """
    Semantic chunking: splits text at sentence boundaries while respecting chunk size.
    Preserves semantic coherence by not breaking mid-sentence.
    """
    # Split into sentences
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = []
    current_size = 0
    
    for sentence in sentences:
        words = sentence.split()
        sentence_size = len(words)
        
        if current_size + sentence_size > chunk_size and current_chunk:
            # Save current chunk and start new one
            chunks.append(' '.join(current_chunk))
            current_chunk = [sentence]
            current_size = sentence_size
        else:
            current_chunk.append(sentence)
            current_size += sentence_size
    
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks

def sliding_window_chunking(text: str, window_size: int = SLIDING_WINDOW_SIZE, 
                           overlap: int = WINDOW_OVERLAP) -> List[str]:
    """
    Sliding window chunking: creates overlapping chunks to maintain context.
    Helps ensure important information isn't lost at chunk boundaries.
    """
    words = text.split()
    chunks = []
    
    for i in range(0, len(words), window_size - overlap):
        chunk = ' '.join(words[i:i + window_size])
        if chunk.strip():
            chunks.append(chunk)
        
        # Stop if we're at the end
        if i + window_size >= len(words):
            break
    
    return chunks

def hybrid_chunking(pages: List[Dict]) -> List[Dict]:
    """
    Hybrid approach: combines semantic and sliding window chunking.
    - Semantic chunks for coherent content
    - Sliding windows to ensure no information is lost at boundaries
    """
    all_chunks = []
    chunk_id = 0
    
    for p in pages:
        text = p["text"]
        
        # First apply semantic chunking
        semantic_chunks = semantic_chunking(text, SEMANTIC_CHUNK_SIZE)
        
        for sem_chunk in semantic_chunks:
            # Add semantic chunk
            all_chunks.append({
                "page": p["page"],
                "text": sem_chunk,
                "doc": p.get("doc", "unknown"),
                "chunk_type": "semantic",
                "chunk_id": chunk_id
            })
            chunk_id += 1
        
        # Also apply sliding window for overlap coverage
        if len(text.split()) > SLIDING_WINDOW_SIZE:
            window_chunks = sliding_window_chunking(text, SLIDING_WINDOW_SIZE, WINDOW_OVERLAP)
            
            for win_chunk in window_chunks:
                all_chunks.append({
                    "page": p["page"],
                    "text": win_chunk,
                    "doc": p.get("doc", "unknown"),
                    "chunk_type": "window",
                    "chunk_id": chunk_id
                })
                chunk_id += 1
    
    return all_chunks

# ================= MULTI-STAGE RETRIEVAL WITH RERANKING =================
def retrieve_and_rerank(question: str, collection, top_k: int = FINAL_CONTEXT_K) -> List[Dict]:
    """
    Multi-stage retrieval:
    1. Initial retrieval: Get top N candidates using semantic similarity
    2. Reranking: Use cross-encoder to rerank candidates based on relevance
    3. Deduplication: Remove similar/duplicate chunks
    4. Return top K most relevant chunks
    """
    # Stage 1: Initial retrieval with higher K
    q_emb = embedding_model.encode([question]).tolist()
    res = collection.query(q_emb, n_results=INITIAL_RETRIEVAL_K)
    
    # Prepare candidates
    candidates = []
    for doc, metadata, distance in zip(res["documents"][0], res["metadatas"][0], res["distances"][0]):
        candidates.append({
            "text": doc,
            "page": metadata.get("page", 1),
            "doc": metadata.get("doc", "unknown"),
            "chunk_type": metadata.get("chunk_type", "unknown"),
            "initial_score": 1 - distance  # Convert distance to similarity
        })
    
    if not candidates:
        return []
    
    # Stage 2: Rerank using cross-encoder
    pairs = [[question, c["text"]] for c in candidates]
    rerank_scores = cross_encoder.predict(pairs)
    
    for i, score in enumerate(rerank_scores):
        candidates[i]["rerank_score"] = float(score)
    
    # Sort by rerank score
    candidates.sort(key=lambda x: x["rerank_score"], reverse=True)
    
    # Stage 3: Deduplication - remove very similar chunks
    final_chunks = []
    seen_texts = []
    
    for candidate in candidates[:RERANK_TOP_K]:
        # Check if this chunk is too similar to already selected ones
        is_duplicate = False
        for seen in seen_texts:
            # Simple similarity check based on text overlap
            overlap = len(set(candidate["text"].split()) & set(seen.split()))
            total = len(set(candidate["text"].split()) | set(seen.split()))
            if total > 0 and overlap / total > 0.7:  # 70% overlap threshold
                is_duplicate = True
                break
        
        if not is_duplicate:
            final_chunks.append(candidate)
            seen_texts.append(candidate["text"])
            
            if len(final_chunks) >= top_k:
                break
    
    return final_chunks

# ================= MULTI-QUESTION PARSING =================
def parse_multiple_questions(user_input: str) -> List[str]:
    questions = re.split(r'\n+|\?+(?=\s*[A-Z0-9])', user_input)
    parsed = []
    
    for q in questions:
        q = q.strip()
        if q and len(q) > 5:
            if not q.endswith('?'):
                q += '?'
            parsed.append(q)
    
    if not parsed:
        parsed = [user_input.strip()]
    
    return parsed

# ================= ENHANCED LLM QUERY WITH BETTER PROMPTING =================
def llm_query(question: str, chunks: List[Dict]) -> str:
    """
    Enhanced LLM query with structured prompting for better answers.
    Instructs the model to provide well-formatted, comprehensive responses.
    """
    if not chunks:
        return "❌ Answer not found in the documents."
    
    # Build rich context with source attribution
    context_parts = []
    for i, c in enumerate(chunks, 1):
        score_info = f"[Relevance: {c.get('rerank_score', 0):.2f}]"
        source_info = f"**Source {i}** (Page {c['page']}, {c.get('doc', 'document')}) {score_info}:"
        context_parts.append(f"{source_info}\n{c['text']}\n")
    
    context = "\n".join(context_parts)
    
    # Enhanced prompt with explicit formatting instructions
    prompt = f"""You are a knowledgeable teaching assistant. Answer the question using ONLY the document content provided below.

**INSTRUCTIONS:**
1. Provide a consice and well-structured answer
2. Use bullet points for listing multiple items or key points
3. Use paragraphs for explanations and descriptions
4. Include specific details, dates, numbers, and examples from the documents
6. Add very short explanations to help understanding, don't add irrelavent information
7. Cite the source number (e.g., "Source 1") when referencing specific information
8. If the answer is not in the documents, clearly state: "Answer not found in the documents."

**DOCUMENT CONTENT:**
{context}

**QUESTION:** {question}

**ANSWER:**"""
    
    payload = {
        "model": MODEL_NAME,
        "messages": [
            {
                "role": "system", 
                "content": "You are an expert teaching assistant who provides clear, well-structured, and comprehensive answers based on document content. You format answers with appropriate structure including bullet points, paragraphs, and explanations for easy understanding."
            },
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.2,
        "max_tokens": 600,
        "top_p": 0.9
    }
    
    try:
        r = requests.post(HF_API_URL, headers=HEADERS, json=payload, timeout=30)
        
        if r.status_code != 200:
            return f"⚠️ API Error: {r.status_code}"
        
        result = r.json()
        answer = result["choices"][0]["message"]["content"].strip()
        
        return answer
    except Exception as e:
        return f"⚠️ Error querying LLM: {str(e)}"

def process_multiple_questions(questions: List[str], collection) -> Dict[str, Dict]:
    """Process multiple questions with enhanced retrieval"""
    results = {}
    
    for q in questions:
        # Use multi-stage retrieval with reranking
        chunks = retrieve_and_rerank(q, collection, FINAL_CONTEXT_K)
        
        # Get answer
        answer = llm_query(q, chunks)
        
        # Store answer with metadata
        results[q] = {
            "answer": answer,
            "chunks_used": len(chunks),
            "sources": [{"doc": c.get("doc"), "page": c.get("page"), "score": c.get("rerank_score", 0)} for c in chunks]
        }
    
    return results

# ================= SESSION STATE =================
if "ready" not in st.session_state:
    st.session_state.ready = False
if "processed_files" not in st.session_state:
    st.session_state.processed_files = []
if "chunk_stats" not in st.session_state:
    st.session_state.chunk_stats = {}

# ================= UI =================
st.header("📤 Upload Documents")

# Settings in sidebar
with st.sidebar:
    st.header("⚙️ Advanced Settings")
    
    st.subheader("Retrieval Parameters")
    initial_k = st.slider("Initial Retrieval K", 10, 50, INITIAL_RETRIEVAL_K, 5,
                         help="Number of chunks to retrieve initially")
    final_k = st.slider("Final Context K", 3, 15, FINAL_CONTEXT_K, 1,
                       help="Number of top chunks to use for answer")
    
    st.subheader("Chunking Parameters")
    semantic_size = st.slider("Semantic Chunk Size (words)", 200, 800, SEMANTIC_CHUNK_SIZE, 50)
    window_size = st.slider("Sliding Window Size (words)", 300, 1000, SLIDING_WINDOW_SIZE, 50)
    window_overlap = st.slider("Window Overlap (words)", 50, 300, WINDOW_OVERLAP, 25)

uploaded_files = st.file_uploader(
    "Upload multiple documents (PDF, DOCX, PPTX, TXT, MD)",
    type=["pdf", "docx", "doc", "pptx", "ppt", "txt", "md"],
    accept_multiple_files=True
)

if uploaded_files and st.button("🚀 Process All Documents with Advanced Chunking"):
    with st.spinner("Processing documents with hybrid chunking strategy..."):
        all_pages = []
        all_images = []
        
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        for idx, uploaded in enumerate(uploaded_files):
            status_text.info(f"📄 Processing: {uploaded.name}")
            
            file_bytes = uploaded.read()
            pages = extract_text_from_file(file_bytes, uploaded.name)
            
            for p in pages:
                p["doc"] = uploaded.name
            
            all_pages.extend(pages)
            
            if uploaded.name.lower().endswith('.pdf'):
                images = extract_images_from_pdf(file_bytes, uploaded.name)
                all_images.extend(images)
            
            progress_bar.progress((idx + 1) / len(uploaded_files))
        
        status_text.info("🔍 Performing OCR on extracted images...")
        for img in all_images:
            text = ocr_image_text(img["path"])
            if text:
                all_pages.append({
                    "page": img["page"],
                    "text": f"[Image OCR] {text}",
                    "doc": img["doc"]
                })
        
        status_text.info("✂️ Applying hybrid chunking (semantic + sliding window)...")
        chunks = hybrid_chunking(all_pages)
        
        if chunks:
            texts = [c["text"] for c in chunks]
            
            status_text.info("🧠 Generating embeddings...")
            embeddings = embedding_model.encode(texts, show_progress_bar=False).tolist()
            
            try:
                chroma_client.delete_collection("syllabus")
            except:
                pass
            
            status_text.info("💾 Storing in vector database...")
            col = chroma_client.create_collection("syllabus")
            col.add(
                documents=texts,
                embeddings=embeddings,
                metadatas=[
                    {
                        "page": c["page"], 
                        "doc": c.get("doc", "unknown"),
                        "chunk_type": c.get("chunk_type", "unknown")
                    } for c in chunks
                ],
                ids=[f"c{i}" for i in range(len(texts))]
            )
            
            # Store statistics
            semantic_count = sum(1 for c in chunks if c.get("chunk_type") == "semantic")
            window_count = sum(1 for c in chunks if c.get("chunk_type") == "window")
            
            st.session_state.chunk_stats = {
                "total": len(chunks),
                "semantic": semantic_count,
                "window": window_count
            }
            st.session_state.ready = True
            st.session_state.processed_files = [f.name for f in uploaded_files]
            
            status_text.empty()
            progress_bar.empty()
            
            st.success("✅ Processing Complete!")
            
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Documents", len(uploaded_files))
            with col2:
                st.metric("Total Chunks", len(chunks))
            with col3:
                st.metric("Semantic Chunks", semantic_count)
            with col4:
                st.metric("Window Chunks", window_count)
            
            st.info(f"📄 **Processed Files:** {', '.join(st.session_state.processed_files)}")
        else:
            st.error("❌ No content extracted from documents")

# ================= QUERY INTERFACE =================
if st.session_state.ready:
    st.divider()
    st.header("💬 Ask Questions")
    
    st.info("💡 **Tips:** You can ask multiple questions (separate with new lines). The system uses multi-stage retrieval with reranking for best results.")
    
    # Show chunk statistics
    with st.expander("📊 Document Statistics", expanded=False):
        stats = st.session_state.chunk_stats
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Total Chunks", stats.get("total", 0))
        with col2:
            st.metric("Semantic Chunks", stats.get("semantic", 0))
        with col3:
            st.metric("Sliding Window Chunks", stats.get("window", 0))
    
    user_input = st.text_area(
        "Enter your question(s)",
        height=150,
        placeholder="Example:\nWhat are the course objectives?\nWhen is the final exam?\nWhat textbook is required?"
    )
    
    if st.button("🔍 Get Answers", type="primary") and user_input:
        questions = parse_multiple_questions(user_input)
        
        st.subheader(f"📝 Processing {len(questions)} question(s)")
        
        col = chroma_client.get_collection("syllabus")
        
        with st.spinner("🔄 Retrieving and reranking relevant content..."):
            results = process_multiple_questions(questions, col)
        
        # Display results with enhanced formatting
        for i, (question, result) in enumerate(results.items(), 1):
            with st.container():
                st.markdown(f"### ❓ Question {i}")
                st.markdown(f"**{question}**")
                
                # Answer section
                st.markdown("#### 📖 Answer")
                answer = result["answer"]
                
                # Format the answer in a nice box
                st.markdown(f"""
                <div style="background-color: #000000; padding: 20px; border-radius: 10px; border-left: 5px solid #4CAF50;">
                {answer}
                </div>
                """, unsafe_allow_html=True)
                
                # Source information in expander
                with st.expander("📚 Sources Used", expanded=False):
                    st.caption(f"**Chunks analyzed:** {result['chunks_used']}")
                    
                    for j, source in enumerate(result['sources'], 1):
                        score_color = "#4CAF50" if source['score'] > 0.5 else "#FF9800"
                        st.markdown(f"""
                        **Source {j}:** {source['doc']} (Page {source['page']})  
                        <span style="color: {score_color}">Relevance Score: {source['score']:.3f}</span>
                        """, unsafe_allow_html=True)
                
                st.divider()

# ================= RESET =================
st.divider()
col1, col2, col3 = st.columns(3)

with col1:
    if st.button("🔄 Reset Database"):
        try:
            chroma_client.delete_collection("syllabus")
        except:
            pass
        shutil.rmtree(IMAGE_DIR, ignore_errors=True)
        os.makedirs(IMAGE_DIR, exist_ok=True)
        st.session_state.ready = False
        st.session_state.processed_files = []
        st.session_state.chunk_stats = {}
        st.success("✅ Reset complete")
        st.rerun()

with col2:
    if st.session_state.ready:
        st.metric("📚 Documents Loaded", len(st.session_state.processed_files))

with col3:
    if st.session_state.ready:
        st.metric("🧩 Total Chunks", st.session_state.chunk_stats.get("total", 0))

# ================= FOOTER =================
st.divider()
st.caption("""
**Advanced Features:**
- 🎯 Multi-stage retrieval with cross-encoder reranking
- ✂️ Hybrid chunking: semantic + sliding window
- 🔍 Deduplication to avoid redundant content
- 📊 Relevance scoring for transparency
- 🎨 Enhanced answer formatting
""")